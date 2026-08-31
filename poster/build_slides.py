"""Build the 6-minute LivSURF presentation deck.

    python build_slides.py      ->  ../video_slides.pptx

Nine slides, 16:9, timed against video_outline.md. The palette is taken from the
figures themselves - the same navy, ochre and green matplotlib now draws with - so a
plotted curve and a slide heading are the same colour rather than two near-misses.
"""
import os

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

HERE = os.path.dirname(os.path.abspath(__file__))
SRC = os.path.abspath(os.path.join(HERE, os.pardir))
OUT = os.path.join(SRC, 'video_slides.pptx')

# from the notebook's BASE_STYLE colour cycle
NAVY = RGBColor(0x1F, 0x3B, 0x63)
OCHRE = RGBColor(0xA8, 0x57, 0x1F)
GREEN = RGBColor(0x2E, 0x6F, 0x55)
INK = RGBColor(0x1B, 0x1A, 0x18)
SOFT = RGBColor(0x57, 0x53, 0x4B)
PAPER = RGBColor(0xFF, 0xFF, 0xFF)
TINT = RGBColor(0xEC, 0xEF, 0xF4)

W, H = Inches(13.333), Inches(7.5)
M = Inches(0.75)                      # side margin
BODY = 'Calibri'
HEAD = 'Cambria'

prs = Presentation()
prs.slide_width, prs.slide_height = W, H
BLANK = prs.slide_layouts[6]


def slide(dark=False):
    s = prs.slides.add_slide(BLANK)
    bg = s.background.fill
    bg.solid()
    bg.fore_color.rgb = NAVY if dark else PAPER
    return s


def text(s, x, y, w, h, runs, size=16, color=INK, font=BODY, bold=False,
         align=PP_ALIGN.LEFT, space=6, line=None):
    """runs: a string, or a list of (text, bold, colour) tuples per paragraph."""
    box = s.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = 0
    tf.margin_top = tf.margin_bottom = 0
    items = [runs] if isinstance(runs, str) else runs
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space)
        if line:
            p.line_spacing = line
        pieces = [(item, bold, color)] if isinstance(item, str) else item
        for t, b, c in pieces:
            r = p.add_run()
            r.text = t
            r.font.size = Pt(size)
            r.font.bold = b
            r.font.name = font
            r.font.color.rgb = c
    return box


def title(s, txt, dark=False, size=38):
    text(s, M, Inches(0.62), W - 2 * M, Inches(1.0), txt, size=size, bold=True,
         font=HEAD, color=PAPER if dark else NAVY)


def stat(s, x, y, value, label, colour=NAVY, vsize=54, w=Inches(3.3)):
    text(s, x, y, w, Inches(0.95), value, size=vsize, bold=True,
         font=HEAD, color=colour, space=0)
    text(s, x, y + Inches(0.1) + Inches(vsize / 72.0), w, Inches(0.8), label,
         size=13, color=SOFT, space=0, line=1.25)


def card(s, x, y, w, h):
    sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    sh.fill.solid()
    sh.fill.fore_color.rgb = TINT
    sh.line.fill.background()
    sh.shadow.inherit = False
    sh.adjustments[0] = 0.06
    return sh


def picture(s, name, y, height):
    """Place a figure centred, scaled to `height`."""
    from PIL import Image
    path = os.path.join(SRC, name)
    iw, ih = Image.open(path).size
    w = Emu(int(height * iw / ih))
    s.shapes.add_picture(path, Emu(int((W - w) / 2)), y, width=w, height=height)


def block(s, x, y, w, n, head, body, colour, tall=False):
    """Numbered step: big numeral, heading, description.

    `tall` leaves room for a heading that wraps to two lines; without it the
    description is placed at a fixed offset and collides with the second line.
    """
    text(s, x, y, Inches(0.7), Inches(0.7), n, size=34, bold=True, font=HEAD,
         color=colour, space=0)
    text(s, x + Inches(0.75), y + Inches(0.05), w - Inches(0.75), Inches(0.8),
         head, size=17, bold=True, color=INK, space=3)
    text(s, x + Inches(0.75), y + Inches(0.92 if tall else 0.52), w - Inches(0.75),
         Inches(1.2), body, size=13.5, color=SOFT, space=0, line=1.25)


# ----------------------------------------------------------------- 1. title
s = slide(dark=True)
text(s, M, Inches(2.15), W - 2 * M, Inches(1.9),
     'AI-Native Transceiver Design\nfor Next-Generation Wireless Networks',
     size=40, bold=True, font=HEAD, color=PAPER, line=1.15)
text(s, M, Inches(4.35), W - 2 * M, Inches(0.5),
     'Yuxing Mao  and  Haikun Xu', size=20, color=RGBColor(0xCA, 0xDC, 0xFC))
text(s, M, Inches(4.95), W - 2 * M, Inches(0.9),
     [[('Supervisor: Dr Huynh Nguyen', False, RGBColor(0xA9, 0xC3, 0xDC))],
      [('LivSURF 436026  ·  School of Computer Science and Informatics', False,
        RGBColor(0xA9, 0xC3, 0xDC))]], size=14, space=3)
s.shapes.add_picture(os.path.join(SRC, 'uol_logo_white.png'), M, Inches(0.7),
                     height=Inches(0.62))
s.notes_slide.notes_text_frame.text = (
    'A, 0:00-0:10. Introduce both of us and the project in one sentence: '
    'using deep learning to redesign a wireless transceiver.')

# ------------------------------------------------------- 2. problem + questions
s = slide()
title(s, 'The problem, and what we asked')
text(s, M, Inches(1.6), Inches(6.0), Inches(2.2),
     [[('A conventional transceiver design treats coding, modulation, channel '
        'estimation and detection as separate blocks.', False, INK)],
      [('Each block is optimal for its own task, but not for the system as a whole. '
        'Consequently, the end-to-end performance is limited.', False, INK)],
      [('Instead, the transmitter and receiver blocks can be replaced by deep neural '
        'networks, jointly optimised as an autoencoder.', False, INK)]],
     size=16, space=10, line=1.3)
c1 = card(s, M, Inches(4.35), Inches(6.0), Inches(1.05))
text(s, M + Inches(0.3), Inches(4.57), Inches(5.4), Inches(0.7),
     'Training needs a differentiable channel model. In practice the channel is a '
     'black box: only its inputs and outputs can be observed.',
     size=14, color=INK, line=1.25)
c2 = card(s, M, Inches(5.6), Inches(6.0), Inches(1.05))
text(s, M + Inches(0.3), Inches(5.82), Inches(5.4), Inches(0.7),
     'Gains are usually shown on AWGN, where the classical receiver is already optimal.',
     size=14, color=INK, line=1.25)
text(s, Inches(7.4), Inches(1.6), Inches(5.2), Inches(0.5), 'Research questions',
     size=20, bold=True, font=HEAD, color=NAVY, space=12)
block(s, Inches(7.4), Inches(2.3), Inches(5.2), '1',
      'How much is learning worth?',
      'Measured against a classical receiver that is already doing its best.', NAVY)
block(s, Inches(7.4), Inches(4.1), Inches(5.2), '2',
      'Can the end-to-end system train with no channel model?',
      'The channel is sampled, never differentiated.', OCHRE, tall=True)
s.notes_slide.notes_text_frame.text = (
    'A, 0:10-0:40. The two objections lead to the two questions. Do not read the '
    'cards out; say them in your own words.')

# ------------------------------------------------------------- 3. approach
s = slide()
title(s, 'Making the comparison fair')
boxes = [('LDPC\nencoder', False), ('Mapper', True), ('Channel', False),
         ('Demapper', True), ('LDPC\ndecoder', False)]
bw, bh, gap = Inches(1.95), Inches(1.15), Inches(0.42)
x0 = M
for i, (lbl, shaded) in enumerate(boxes):
    x = x0 + i * (bw + gap)
    sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(1.65), bw, bh)
    sh.fill.solid()
    sh.fill.fore_color.rgb = TINT if shaded else PAPER
    sh.line.color.rgb = NAVY
    sh.line.width = Pt(1.25)
    sh.shadow.inherit = False
    sh.adjustments[0] = 0.10
    tf = sh.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = lbl
    r.font.size = Pt(15)
    r.font.bold = shaded
    r.font.name = BODY
    r.font.color.rgb = NAVY if shaded else INK
    if i < len(boxes) - 1:
        con = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, x + bw + Inches(0.09),
                                 Inches(2.03), Inches(0.24), Inches(0.24))
        con.fill.solid()
        con.fill.fore_color.rgb = NAVY
        con.line.fill.background()
        con.shadow.inherit = False
text(s, M, Inches(3.0), W - 2 * M, Inches(0.4),
     'Only the two shaded blocks differ. Same LDPC code, same channel, same random '
     'seeds, and both systems evaluated on the same terms.',
     size=14, color=SOFT)
text(s, M, Inches(3.9), Inches(5.9), Inches(0.5), 'Pilots are not free',
     size=19, bold=True, font=HEAD, color=NAVY, space=8)
text(s, M, Inches(4.5), Inches(5.9), Inches(1.9),
     [[('They occupy symbol slots, so the effective code rate falls to R(L-P)/L.',
        False, INK)],
      [('The channel estimate is imperfect, which adds noise: a factor of (1+1/P).',
        False, INK)]],
     size=15, space=8, line=1.3)
stat(s, Inches(7.4), Inches(3.85), '1.76 dB', 'what that factor is worth at two pilots. '
     'Before we added it, we were understating our own noise.', OCHRE, vsize=46)
stat(s, Inches(7.4), Inches(5.75), '64,000', 'codewords per SNR point. Points with fewer '
     'than 30 block errors are dropped, not plotted.', NAVY, vsize=46)
s.notes_slide.notes_text_frame.text = (
    'A, 0:40-2:00. The heart of your section. Land 1.76 dB and say the sentence about '
    'understating our own noise - it sets up your skill at 3:20.')

# ---------------------------------------------------------- 4. AWGN result
s = slide()
title(s, 'On AWGN, learning wins, and the transmitter is why')
picture(s, 'awgn_ae_vs_qam_poster.png', Inches(1.65), Inches(3.25))
text(s, M, Inches(5.35), Inches(5.8), Inches(1.6),
     [[('A posteriori probability (APP) demapping is already optimal here, so the '
        'receiver cannot be the source of the gain.', False, INK)],
      [('It comes from the transmitter. The learned constellation leaves the regular '
        'lattice and spends more power on the symbols that are most often confused.',
        False, INK)]],
     size=15, space=8, line=1.3)
stat(s, Inches(9.3), Inches(5.5), '0.37 dB', 'gain over classical 64-QAM at BER 1e-2',
     GREEN, vsize=50)
s.notes_slide.notes_text_frame.text = 'B, 2:00-2:30.'

# ------------------------------------------------------- 5. both channels
s = slide()
title(s, 'Under fading, the gain is small')
rows = [('System', '1e-2', '1e-3', '1e-4'),
        ('AWGN, classical 64-QAM + APP', '5.67', '6.17', '6.54'),
        ('AWGN, learned (channel gradient)', '5.30', '5.78', '6.18'),
        ('AWGN, learned (RL, no channel model)', '5.42', '5.92', '6.31'),
        ('Rayleigh, classical 64-QAM + APP', '11.26', '13.30', '14.79'),
        ('Rayleigh, learned', '11.25', '13.19', '14.94')]
tw, th = Inches(7.6), Inches(3.5)
gt = s.shapes.add_table(len(rows), 4, M, Inches(1.7), tw, th).table
gt.columns[0].width = Inches(3.7)
for c in range(1, 4):
    gt.columns[c].width = Inches(1.3)
for r, row in enumerate(rows):
    gt.rows[r].height = Inches(0.5)
    for c, val in enumerate(row):
        cell = gt.cell(r, c)
        cell.text = val
        cell.margin_left = cell.margin_right = Inches(0.08)
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY if r == 0 else (
            TINT if r == 2 else PAPER)
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT if c == 0 else PP_ALIGN.RIGHT
        f = p.runs[0].font
        f.size = Pt(13)
        f.name = BODY
        f.bold = (r == 0 or r == 2)
        f.color.rgb = PAPER if r == 0 else INK
text(s, M, Inches(5.45), Inches(7.6), Inches(0.4),
     'Eb/N0 needed to reach a given bit error rate, in dB. Lower is better.',
     size=12, color=SOFT)
stat(s, Inches(9.0), Inches(1.75), '5.9 to 8.8 dB', 'the cost of fading', OCHRE, vsize=34, w=Inches(3.8))
stat(s, Inches(9.0), Inches(3.35), 'under 0.4 dB', 'what changing the receiver is worth',
     NAVY, vsize=34, w=Inches(3.8))
text(s, Inches(9.0), Inches(4.95), Inches(3.8), Inches(1.8),
     'Under Rayleigh with estimated CSI there is a small gain at 1e-2 and 1e-3, of '
     '0.01 and 0.11 dB, and a loss of 0.15 dB at 1e-4.',
     size=13.5, color=SOFT, line=1.3)
s.notes_slide.notes_text_frame.text = (
    'B, 2:30-3:00. The point is the vertical gap between the AWGN block and the '
    'Rayleigh block, not the differences inside either one.')

# ------------------------------------------------- 6. no channel model (key)
s = slide()
title(s, 'Training with no channel model at all')
picture(s, 'rl_vs_gradient_awgn.png', Inches(1.6), Inches(3.15))
text(s, M, Inches(5.15), Inches(6.6), Inches(1.7),
     [[('There is no channel layer connecting the two neural networks, because the '
        'channel is not differentiable.', False, INK)],
      [('We feed the training loss from the receiver back to the transmitter, which '
        'is then trained by reinforcement learning. The two networks are still '
        'jointly optimised, without a channel model.', False, INK)]],
     size=15, space=8, line=1.3)
stat(s, Inches(8.3), Inches(5.05), '0.12 to 0.14 dB',
     'what dropping the channel model costs', OCHRE, vsize=28, w=Inches(4.3))
stat(s, Inches(8.3), Inches(6.15), '0.23 to 0.25 dB',
     'still ahead of classical 64-QAM', GREEN, vsize=28, w=Inches(4.3))
s.notes_slide.notes_text_frame.text = (
    'B, 3:00-3:20. This is the result the project turns on. Say both numbers.')

# ------------------------------------------------------------ 7. skill 1
s = slide()
title(s, 'Skill 1  ·  Research rigour')
text(s, M, Inches(1.55), Inches(6.2), Inches(0.5),
     'Checking the thing that would embarrass you', size=19, bold=True,
     font=HEAD, color=OCHRE, space=10)
block(s, M, Inches(2.35), Inches(6.2), '1', 'What happened',
      'Our effective-noise model was missing the (1+1/P) term. Adding it made our own '
      'numbers worse.', NAVY)
block(s, M, Inches(4.55), Inches(6.2), '2', 'It inverted our own conclusion',
      'Two pilots went from apparently best to measurably worst. We changed it anyway, '
      'and wrote a validation cell that tests the model against the simulator.', NAVY)
card(s, Inches(7.5), Inches(2.35), Inches(5.1), Inches(4.2))
text(s, Inches(7.85), Inches(2.65), Inches(4.4), Inches(0.4), 'Beyond LivSURF',
     size=17, bold=True, font=HEAD, color=NAVY, space=10)
text(s, Inches(7.85), Inches(3.3), Inches(4.4), Inches(3.0),
     [[('In research and in industry you will get results that favour you.', False, INK)],
      [('What makes a result trustworthy is whether you were willing to test the thing '
        'that would embarrass you.', False, INK)],
      [('Directly useful for postgraduate study, and for any role that works with data.',
        False, INK)]],
     size=14, space=9, line=1.3)
s.notes_slide.notes_text_frame.text = (
    'A, 3:20-4:20. Tell it as a story. The "Beyond LivSURF" panel is the half the '
    'marking scheme asks for, so do not skip it.')

# ------------------------------------------------------------ 8. skill 2
s = slide()
title(s, 'Skill 2  ·  Problem solving under uncertainty')
text(s, M, Inches(1.55), Inches(6.2), Inches(0.5),
     'Diagnose by intervention, not by argument', size=19, bold=True,
     font=HEAD, color=GREEN, space=10)
block(s, M, Inches(2.35), Inches(6.2), '1', 'A curve that could not be real',
      'Reused on a fading channel, the BER rose as the SNR rose. Our first explanation '
      'blamed deep fades.', NAVY)
block(s, M, Inches(4.55), Inches(6.2), '2', 'We tested instead of arguing',
      'Weights held fixed, one thing changed at a time. Erasing deep fades made it '
      'worse. The real cause was the quiet tail: at 18 dB, 89% of symbols were quieter '
      'than anything seen in training.', NAVY)
card(s, Inches(7.5), Inches(2.35), Inches(5.1), Inches(4.2))
text(s, Inches(7.85), Inches(2.65), Inches(4.4), Inches(0.4), 'Beyond LivSURF',
     size=17, bold=True, font=HEAD, color=NAVY, space=10)
text(s, Inches(7.85), Inches(3.3), Inches(4.4), Inches(3.0),
     [[('The first explanation for a bug is usually the most familiar one, not the '
        'right one.', False, INK)],
      [('Designing an experiment that can prove yourself wrong is faster than arguing '
        'about it.', False, INK)],
      [('That habit transfers to any engineering or analysis work.', False, INK)]],
     size=14, space=9, line=1.3)
s.notes_slide.notes_text_frame.text = 'B, 4:20-5:20. Same shape as slide 7.'

# --------------------------------------------------------- 9. reflection
s = slide(dark=True)
title(s, 'What we take away', dark=True, size=36)
text(s, M, Inches(1.9), Inches(5.7), Inches(2.6),
     [[('Under fading the measured gain is small, and at 1e-4 it goes the other '
        'way.', False, PAPER)],
      [('That looks like a disappointing result, until you notice the interesting '
        'finding was somewhere else entirely.', False, RGBColor(0xCA, 0xDC, 0xFC))]],
     size=18, space=12, line=1.3)
card_bg = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.9), Inches(1.85),
                             Inches(5.7), Inches(3.5))
card_bg.fill.solid()
card_bg.fill.fore_color.rgb = RGBColor(0x2C, 0x4E, 0x7A)
card_bg.line.fill.background()
card_bg.shadow.inherit = False
card_bg.adjustments[0] = 0.06
text(s, Inches(7.3), Inches(2.25), Inches(4.9), Inches(2.8),
     [[('A measurement that shows little is worth as much as one that shows a lot, '
        'provided it is clean enough to trust.', False, PAPER)],
      [('The end-to-end system can be trained with no channel model, and we can now '
        'say what that costs.', False,
        RGBColor(0xCA, 0xDC, 0xFC))]],
     size=16, space=14, line=1.35)
text(s, M, Inches(5.9), Inches(11.8), Inches(0.6),
     'Notebook, figures and measured data: github.com/MaoYuxingy',
     size=13, color=RGBColor(0xA9, 0xC3, 0xDC))
s.notes_slide.notes_text_frame.text = (
    'A and B, 5:20-6:00. A takes the left, B takes the right, then thank the audience.')

prs.save(OUT)
print('%s  %.2f MB, %d slides' % (os.path.basename(OUT),
                                  os.path.getsize(OUT) / 1024 / 1024,
                                  len(prs.slides._sldIdLst)))
